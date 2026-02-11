import streamlit as st
import requests
import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import time
import json
import pandas as pd
import zipfile
import matplotlib.pyplot as plt
from io import BytesIO
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from datetime import datetime
import hashlib
import sqlite3

# ============================================================================
# PAGE CONFIGURATION
# ============================================================================

st.set_page_config(
    page_title="AI PowerPoint Generator Pro",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ============================================================================
# DATABASE FUNCTIONS
# ============================================================================

def init_database():
    """Initialize database with migration support"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    
    # Create users table
    c.execute('''CREATE TABLE IF NOT EXISTS users
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  username TEXT UNIQUE NOT NULL,
                  password_hash TEXT NOT NULL,
                  email TEXT,
                  created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  last_login TIMESTAMP,
                  is_active BOOLEAN DEFAULT 1,
                  role TEXT DEFAULT 'user')''')
    
    # Create usage_logs table
    c.execute('''CREATE TABLE IF NOT EXISTS usage_logs
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id INTEGER,
                  action TEXT,
                  topic TEXT,
                  slides_count INTEGER,
                  timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  FOREIGN KEY (user_id) REFERENCES users (id))''')
    
    # Create sessions table
    c.execute('''CREATE TABLE IF NOT EXISTS sessions
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id INTEGER,
                  login_time TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                  logout_time TIMESTAMP,
                  is_active BOOLEAN DEFAULT 1,
                  session_token TEXT,
                  FOREIGN KEY (user_id) REFERENCES users (id))''')
    
    # Migration: Add missing columns
    try:
        c.execute("PRAGMA table_info(sessions)")
        columns = [column[1] for column in c.fetchall()]
        
        if 'is_active' not in columns:
            c.execute("ALTER TABLE sessions ADD COLUMN is_active BOOLEAN DEFAULT 1")
            conn.commit()
            
        if 'session_token' not in columns:
            c.execute("ALTER TABLE sessions ADD COLUMN session_token TEXT")
            conn.commit()
    except Exception as e:
        pass
    
    # Create admin user if not exists
    c.execute("SELECT * FROM users WHERE username = 'admin'")
    if not c.fetchone():
        admin_password = hashlib.sha256('admin123'.encode()).hexdigest()
        c.execute("INSERT INTO users (username, password_hash, email, role) VALUES (?, ?, ?, ?)",
                  ('admin', admin_password, 'admin@pptgen.com', 'admin'))
    
    conn.commit()
    conn.close()

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def verify_user(username, password):
    """Verify and login user"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    
    try:
        password_hash = hash_password(password)
        c.execute("SELECT id, username, role, is_active FROM users WHERE username = ? AND password_hash = ?",
                  (username, password_hash))
        
        user = c.fetchone()
        
        if user and user[3]:
            c.execute("UPDATE users SET last_login = ? WHERE id = ?", (datetime.now(), user[0]))
            session_token = hashlib.md5(f"{user[0]}{datetime.now()}".encode()).hexdigest()
            
            try:
                c.execute("UPDATE sessions SET is_active = 0, logout_time = ? WHERE user_id = ? AND is_active = 1", 
                          (datetime.now(), user[0]))
            except sqlite3.OperationalError:
                pass
            
            try:
                c.execute("INSERT INTO sessions (user_id, login_time, is_active, session_token) VALUES (?, ?, ?, ?)",
                          (user[0], datetime.now(), 1, session_token))
            except sqlite3.OperationalError:
                c.execute("INSERT INTO sessions (user_id, login_time) VALUES (?, ?)",
                          (user[0], datetime.now()))
            
            conn.commit()
            conn.close()
            
            return {
                'id': user[0], 
                'username': user[1], 
                'role': user[2], 
                'is_active': user[3], 
                'session_token': session_token
            }
        
        conn.close()
        return None
    except Exception as e:
        conn.close()
        return None

def create_user_by_admin(username, password, email):
    """Admin creates user"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    try:
        password_hash = hash_password(password)
        c.execute("INSERT INTO users (username, password_hash, email, role) VALUES (?, ?, ?, ?)",
                  (username, password_hash, email, 'user'))
        conn.commit()
        conn.close()
        return True
    except sqlite3.IntegrityError:
        conn.close()
        return False

def logout_user(user_id):
    """Logout user"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    try:
        c.execute("UPDATE sessions SET is_active = 0, logout_time = ? WHERE user_id = ? AND is_active = 1",
                  (datetime.now(), user_id))
    except sqlite3.OperationalError:
        c.execute("UPDATE sessions SET logout_time = ? WHERE user_id = ? AND logout_time IS NULL",
                  (datetime.now(), user_id))
    conn.commit()
    conn.close()

def log_usage(user_id, action, topic="", slides_count=0):
    """Log activity"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("INSERT INTO usage_logs (user_id, action, topic, slides_count) VALUES (?, ?, ?, ?)",
              (user_id, action, topic, slides_count))
    conn.commit()
    conn.close()

def get_user_stats(user_id):
    """Get user stats"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT COUNT(*) FROM usage_logs WHERE user_id = ? AND action = 'generate_presentation'", (user_id,))
    total_presentations = c.fetchone()[0]
    c.execute("SELECT SUM(slides_count) FROM usage_logs WHERE user_id = ? AND action = 'generate_presentation'", (user_id,))
    total_slides = c.fetchone()[0] or 0
    c.execute("SELECT COUNT(*) FROM sessions WHERE user_id = ?", (user_id,))
    total_logins = c.fetchone()[0]
    conn.close()
    return {'total_presentations': total_presentations, 'total_slides': total_slides, 'total_logins': total_logins}

def get_all_users():
    """Get all users"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT id, username, email, created_at, last_login, is_active, role FROM users ORDER BY created_at DESC")
    users = c.fetchall()
    conn.close()
    return users

def get_currently_logged_in_users():
    """Get currently logged in users"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    try:
        c.execute("""
            SELECT u.id, u.username, u.email, s.login_time, u.role
            FROM sessions s
            JOIN users u ON s.user_id = u.id
            WHERE s.is_active = 1
            ORDER BY s.login_time DESC
        """)
        active_users = c.fetchall()
    except sqlite3.OperationalError:
        active_users = []
    conn.close()
    return active_users

def get_user_activity_details(user_id):
    """Get detailed activity for a specific user"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""
        SELECT action, topic, slides_count, timestamp
        FROM usage_logs
        WHERE user_id = ?
        ORDER BY timestamp DESC
        LIMIT 20
    """, (user_id,))
    activities = c.fetchall()
    conn.close()
    return activities

def get_all_user_activities():
    """Get all activities from all users"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("""
        SELECT u.username, l.action, l.topic, l.slides_count, l.timestamp
        FROM usage_logs l
        JOIN users u ON l.user_id = u.id
        ORDER BY l.timestamp DESC
        LIMIT 100
    """)
    activities = c.fetchall()
    conn.close()
    return activities

def get_system_stats():
    """Get system stats"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("SELECT COUNT(*) FROM users WHERE role = 'user'")
    total_users = c.fetchone()[0]
    try:
        c.execute("SELECT COUNT(*) FROM sessions WHERE is_active = 1")
        currently_online = c.fetchone()[0]
    except sqlite3.OperationalError:
        currently_online = 0
    c.execute("SELECT COUNT(*) FROM usage_logs WHERE action = 'generate_presentation'")
    total_presentations = c.fetchone()[0]
    c.execute("SELECT SUM(slides_count) FROM usage_logs WHERE action = 'generate_presentation'")
    total_slides = c.fetchone()[0] or 0
    c.execute("SELECT COUNT(*) FROM sessions WHERE DATE(login_time) = DATE('now')")
    today_logins = c.fetchone()[0]
    conn.close()
    return {
        'total_users': total_users,
        'currently_online': currently_online,
        'total_presentations': total_presentations,
        'total_slides': total_slides,
        'today_logins': today_logins
    }

def toggle_user_status(user_id, is_active):
    """Enable/disable user"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("UPDATE users SET is_active = ? WHERE id = ?", (is_active, user_id))
    conn.commit()
    conn.close()

def delete_user(user_id):
    """Delete user"""
    conn = sqlite3.connect('ppt_generator.db', check_same_thread=False)
    c = conn.cursor()
    c.execute("DELETE FROM users WHERE id = ?", (user_id,))
    conn.commit()
    conn.close()

# ============================================================================
# TEMPLATE MANAGEMENT FUNCTIONS
# ============================================================================

def generate_template_id():
    """Generate unique template ID"""
    return hashlib.md5(str(datetime.now().timestamp()).encode()).hexdigest()[:8]

def save_template_to_state(name, template_data):
    """Save template to session state"""
    template_id = generate_template_id()
    template_data['id'] = template_id
    template_data['name'] = name
    template_data['created_at'] = datetime.now().strftime("%Y-%m-%d %H:%M")
    template_data['usage_count'] = 0
    st.session_state.templates[template_id] = template_data
    return template_id

def load_template_from_state(template_id):
    """Load template from session state"""
    return st.session_state.templates.get(template_id, None)

def delete_template(template_id):
    """Delete template from session state"""
    if template_id in st.session_state.templates:
        del st.session_state.templates[template_id]
        return True
    return False

def export_all_templates():
    """Export all templates as JSON"""
    return json.dumps(st.session_state.templates, indent=2)

def import_templates(json_data):
    """Import templates from JSON"""
    try:
        templates = json.loads(json_data)
        st.session_state.templates.update(templates)
        return True
    except:
        return False

def get_preset_templates():
    """Get preset professional templates"""
    return {
        "pitch_deck": {
            "name": "🚀 Startup Pitch Deck",
            "category": "Pitch",
            "slide_count": 10,
            "tone": "Persuasive",
            "audience": "Investors",
            "theme": "Gradient Modern",
            "image_mode": "With Images",
            "language": "English",
            "description": "Perfect for startup pitches with 10-slide structure"
        },
        "corporate_report": {
            "name": "📈 Corporate Report",
            "category": "Business",
            "slide_count": 12,
            "tone": "Formal",
            "audience": "Corporate",
            "theme": "Corporate Blue",
            "image_mode": "With Images",
            "language": "English",
            "description": "Professional business reporting format"
        },
        "training_session": {
            "name": "🎓 Training Session",
            "category": "Training",
            "slide_count": 15,
            "tone": "Educational",
            "audience": "Students",
            "theme": "Pastel Soft",
            "image_mode": "With Images",
            "language": "English",
            "description": "Educational content with clear structure"
        },
        "sales_pitch": {
            "name": "💼 Sales Pitch",
            "category": "Sales",
            "slide_count": 8,
            "tone": "Persuasive",
            "audience": "Clients",
            "theme": "Professional Green",
            "image_mode": "With Images",
            "language": "English",
            "description": "Compelling sales presentation format"
        },
        "tech_overview": {
            "name": "🔧 Technical Overview",
            "category": "Technical",
            "slide_count": 10,
            "tone": "Neutral",
            "audience": "Managers",
            "theme": "Minimal Dark",
            "image_mode": "With Images",
            "language": "English",
            "description": "Technical documentation and overview"
        },
        "marketing_campaign": {
            "name": "📣 Marketing Campaign",
            "category": "Marketing",
            "slide_count": 9,
            "tone": "Inspirational",
            "audience": "Corporate",
            "theme": "Elegant Purple",
            "image_mode": "With Images",
            "language": "English",
            "description": "Creative marketing strategy presentation"
        }
    }

# ============================================================================
# IMAGE GENERATION FUNCTIONS
# ============================================================================

def generate_topic_search_terms(main_topic, slide_title, image_prompt):
    """Generate search terms prioritizing topic relevance"""
    search_terms = []
    
    if image_prompt and image_prompt.strip():
        search_terms.append(image_prompt.strip())
    
    if main_topic and slide_title:
        search_terms.append(f"{main_topic} {slide_title}")
    
    if slide_title:
        search_terms.append(slide_title)
    
    if main_topic:
        search_terms.append(main_topic)
    
    seen = set()
    unique = []
    for term in search_terms:
        lower = term.lower().strip()
        if lower and lower not in seen:
            seen.add(lower)
            unique.append(term)
    
    return unique

def get_google_image(query, api_key, cx):
    """Get image using Google Custom Search API"""
    try:
        st.session_state.google_searches_used += 1
        
        url = "https://www.googleapis.com/customsearch/v1"
        
        params = {
            'key': api_key,
            'cx': cx,
            'q': query,
            'searchType': 'image',
            'num': 3,
            'imgSize': 'large',
            'imgType': 'photo',
            'safe': 'active',
            'fileType': 'jpg,png'
        }
        
        response = requests.get(url, params=params, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            
            if 'items' in data and len(data['items']) > 0:
                for item in data['items'][:3]:
                    try:
                        image_url = item['link']
                        img_response = requests.get(image_url, timeout=10, headers={
                            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
                        })
                        
                        if img_response.status_code == 200 and len(img_response.content) > 5000:
                            img = Image.open(io.BytesIO(img_response.content))
                            if img.size[0] > 300 and img.size[1] > 200:
                                return img_response.content
                    except:
                        continue
        
        return None
    except Exception as e:
        return None

def get_unsplash_image(query, width=800, height=600):
    """Get image from Unsplash Direct (Free)"""
    try:
        clean_query = query.strip().replace(' ', ',')
        url = f"https://source.unsplash.com/{width}x{height}/?{clean_query}"
        
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        
        response = requests.get(url, timeout=15, allow_redirects=True, headers=headers)
        
        if response.status_code == 200 and len(response.content) > 5000:
            try:
                img = Image.open(io.BytesIO(response.content))
                if img.size[0] > 400 and img.size[1] > 300:
                    return response.content
            except:
                pass
        return None
    except:
        return None

def get_pexels_image(query, api_key):
    """Get image from Pexels Direct API"""
    if not api_key:
        return None
    
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": api_key}
        params = {
            "query": query,
            "per_page": 3,
            "orientation": "landscape"
        }
        
        response = requests.get(url, headers=headers, params=params, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            if data.get("photos"):
                photo = data["photos"][0]
                img_url = photo["src"]["large"]
                
                img_response = requests.get(img_url, timeout=10)
                if img_response.status_code == 200:
                    return img_response.content
        return None
    except:
        return None

def get_topic_relevant_image(main_topic, slide_title, image_prompt, google_api_key, google_cx, use_unsplash, use_pexels, pexels_key):
    """Get highly relevant image using Google + fallbacks"""
    
    search_terms = generate_topic_search_terms(main_topic, slide_title, image_prompt)
    
    for i, term in enumerate(search_terms, 1):
        if google_api_key and google_cx:
            image_data = get_google_image(term, google_api_key, google_cx)
            if image_data:
                return image_data
        
        if use_pexels and pexels_key:
            image_data = get_pexels_image(term, pexels_key)
            if image_data:
                return image_data
        
        if use_unsplash:
            image_data = get_unsplash_image(term)
            if image_data:
                return image_data
    
    fallback = main_topic.split()[0] if main_topic else "business"
    
    if google_api_key and google_cx:
        image_data = get_google_image(fallback, google_api_key, google_cx)
        if image_data:
            return image_data
    
    if use_unsplash:
        image_data = get_unsplash_image(fallback)
        if image_data:
            return image_data
    
    return None

# ============================================================================
# AI CONTENT GENERATION
# ============================================================================

def repair_truncated_json(json_text):
    """Attempt to repair truncated JSON from AI response"""
    text = json_text.strip()
    
    if text.startswith("```json"):
        text = text[7:]
    if text.startswith("```"):
        text = text[3:]
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()
    
    try:
        data = json.loads(text)
        return data
    except json.JSONDecodeError:
        pass
    
    slides = []
    slides_start = text.find('"slides"')
    if slides_start == -1:
        return None
    
    bracket_pos = text.find('[', slides_start)
    if bracket_pos == -1:
        return None
    
    current_pos = bracket_pos + 1
    brace_count = 0
    slide_start = -1
    
    while current_pos < len(text):
        char = text[current_pos]
        
        if char == '{' and brace_count == 0:
            slide_start = current_pos
            brace_count = 1
        elif char == '{':
            brace_count += 1
        elif char == '}':
            brace_count -= 1
            if brace_count == 0 and slide_start != -1:
                slide_text = text[slide_start:current_pos + 1]
                try:
                    slide_obj = json.loads(slide_text)
                    if 'title' in slide_obj:
                        if 'bullets' not in slide_obj:
                            slide_obj['bullets'] = []
                        if 'image_prompt' not in slide_obj:
                            slide_obj['image_prompt'] = slide_obj['title']
                        if 'speaker_notes' not in slide_obj:
                            slide_obj['speaker_notes'] = ""
                        slides.append(slide_obj)
                except:
                    pass
                slide_start = -1
        
        current_pos += 1
    
    if slides:
        return {"slides": slides}
    
    return None

def generate_content_with_claude(api_key, topic, category, slide_count, tone, audience, key_points, model_choice, language, grok_api_key=None, groq_api_key=None):
    """Generate presentation content using AI"""
    try:
        use_grok_api = "Grok" in model_choice and grok_api_key
        use_groq_api = "Groq" in model_choice and groq_api_key
        
        if use_groq_api:
            if "Llama 3.3" in model_choice:
                model = "llama-3.3-70b-versatile"
            else:
                model = "mixtral-8x7b-32768"
            api_url = "https://api.groq.com/openai/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {groq_api_key.strip()}",
                "Content-Type": "application/json",
            }
        elif use_grok_api:
            if "Grok-4" in model_choice:
                model = "grok-4-latest"
            elif "Grok-3" in model_choice:
                model = "grok-3-latest"
            else:
                model = "grok-2-latest"
            api_url = "https://api.x.ai/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {grok_api_key.strip()}",
                "Content-Type": "application/json",
            }
        else:
            if "Gemini" in model_choice:
                model = "google/gemini-2.0-flash-exp:free"
            elif "Llama" in model_choice:
                model = "meta-llama/llama-3.2-3b-instruct:free"
            elif "Mistral" in model_choice:
                model = "mistralai/mistral-7b-instruct:free"
            else:
                model = "anthropic/claude-3.5-sonnet"
            api_url = "https://openrouter.ai/api/v1/chat/completions"
            headers = {
                "Authorization": f"Bearer {api_key.strip()}",
                "Content-Type": "application/json",
            }
        
        calculated_tokens = min(slide_count * 350 + 500, 4000)
        language_instruction = f"Generate ALL content in {language} language." if language != "English" else ""
        
        prompt = f"""{language_instruction}
Create a {slide_count}-slide presentation about: {topic}

Category: {category} | Tone: {tone} | Audience: {audience}
{f"Include these points: {key_points}" if key_points else ""}

Return ONLY valid JSON (no markdown, no extra text):
{{"slides": [
  {{
    "title": "Main Title of Presentation",
    "bullets": [],
    "image_prompt": "professional {topic} banner",
    "speaker_notes": "Welcome and introduction"
  }},
  {{
    "title": "Key Point Title",
    "bullets": ["First important point about the topic", "Second key insight or fact", "Third supporting detail", "Fourth actionable item"],
    "image_prompt": "{topic} concept",
    "speaker_notes": "Explain these points in detail"
  }}
]}}

CRITICAL REQUIREMENTS:
1. First slide is TITLE ONLY (empty bullets array)
2. ALL OTHER SLIDES MUST have 3-5 bullet points
3. Each bullet must be a complete, informative sentence (8-15 words)
4. Total: exactly {slide_count} slides
5. Return ONLY the JSON object

Generate {slide_count} slides now:"""

        response = requests.post(
            api_url,
            headers=headers,
            json={
                "model": model,
                "max_tokens": calculated_tokens,
                "messages": [{"role": "user", "content": prompt}]
            },
            timeout=60
        )
        
        if response.status_code == 200:
            data = response.json()
            content_text = data["choices"][0]["message"]["content"]
            slides_data = repair_truncated_json(content_text)
            
            if slides_data and "slides" in slides_data:
                slides = slides_data["slides"]
                
                if not slides:
                    return None
                
                for i, slide in enumerate(slides):
                    if 'bullets' not in slide:
                        slide['bullets'] = []
                    if 'image_prompt' not in slide:
                        slide['image_prompt'] = slide.get('title', topic)
                    if 'speaker_notes' not in slide:
                        slide['speaker_notes'] = ""
                
                return slides
        return None
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None

def generate_content_with_retry(api_key, topic, category, slide_count, tone, audience, key_points, model_choice, language, grok_api_key=None, groq_api_key=None, max_retries=3):
    """Generate content with automatic retry"""
    for attempt in range(max_retries):
        try:
            result = generate_content_with_claude(api_key, topic, category, slide_count, tone, audience, key_points, model_choice, language, grok_api_key, groq_api_key)
            if result:
                return result
        except Exception as e:
            if "Rate limit" in str(e):
                if attempt < max_retries - 1:
                    wait_time = (attempt + 1) * 5
                    st.warning(f"⏳ Retrying in {wait_time}s...")
                    time.sleep(wait_time)
    return None

# ============================================================================
# ANALYSIS FUNCTIONS
# ============================================================================

def analyze_presentation(slides_content):
    """Analyze presentation quality"""
    issues = []
    suggestions = []
    score = 100
    
    for i, slide in enumerate(slides_content, 1):
        bullet_count = len(slide.get('bullets', []))
        if bullet_count > 5:
            issues.append(f"Slide {i}: Too many bullets ({bullet_count})")
            suggestions.append(f"Slide {i}: Reduce to 3-5 key points")
            score -= 5
        
        title_len = len(slide['title'])
        if title_len > 60:
            issues.append(f"Slide {i}: Title too long ({title_len} chars)")
            suggestions.append(f"Slide {i}: Shorten title to <60 characters")
            score -= 3
    
    score = max(0, score)
    return issues, suggestions, score

# ============================================================================
# EXPORT FUNCTIONS
# ============================================================================

def export_to_pdf(slides_content, topic):
    """Export to PDF"""
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []
    
    title = Paragraph(f"<b>{topic}</b>", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 12))
    
    for i, slide in enumerate(slides_content, 1):
        slide_title = Paragraph(f"<b>Slide {i}: {slide['title']}</b>", styles['Heading2'])
        story.append(slide_title)
        story.append(Spacer(1, 6))
        
        for bullet in slide.get('bullets', []):
            bullet_text = Paragraph(f"• {bullet}", styles['BodyText'])
            story.append(bullet_text)
        
        story.append(Spacer(1, 20))
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def export_to_google_slides_json(slides_content, topic, theme):
    """Export to Google Slides JSON"""
    google_slides_data = {
        "title": topic,
        "theme": theme,
        "slides": []
    }
    
    for slide in slides_content:
        google_slide = {
            "title": slide['title'],
            "content": slide.get('bullets', []),
            "notes": slide.get('speaker_notes', ''),
            "imagePrompt": slide.get('image_prompt', '')
        }
        google_slides_data['slides'].append(google_slide)
    
    return json.dumps(google_slides_data, indent=2)

# ============================================================================
# POWERPOINT CREATION
# ============================================================================

def create_powerpoint(slides_content, theme, image_mode, google_api_key, google_cx, use_unsplash, use_pexels, pexels_key, category, audience, topic, image_position, logo_data, show_progress=True):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    themes = {
        "Corporate Blue": {"bg": RGBColor(240, 248, 255), "accent": RGBColor(31, 119, 180), "text": RGBColor(0, 0, 0)},
        "Gradient Modern": {"bg": RGBColor(240, 242, 246), "accent": RGBColor(138, 43, 226), "text": RGBColor(0, 0, 0)},
        "Minimal Dark": {"bg": RGBColor(30, 30, 30), "accent": RGBColor(255, 215, 0), "text": RGBColor(255, 255, 255)},
        "Pastel Soft": {"bg": RGBColor(255, 250, 240), "accent": RGBColor(255, 182, 193), "text": RGBColor(60, 60, 60)},
        "Professional Green": {"bg": RGBColor(245, 255, 250), "accent": RGBColor(34, 139, 34), "text": RGBColor(0, 0, 0)},
        "Elegant Purple": {"bg": RGBColor(250, 245, 255), "accent": RGBColor(128, 0, 128), "text": RGBColor(0, 0, 0)}
    }
    
    color_scheme = themes.get(theme, themes["Corporate Blue"])
    
    positions = {
        "Right Side": {"left": Inches(6.5), "top": Inches(2), "width": Inches(3)},
        "Left Side": {"left": Inches(0.5), "top": Inches(2), "width": Inches(3)},
        "Top Right Corner": {"left": Inches(8), "top": Inches(0.5), "width": Inches(1.5)},
        "Bottom": {"left": Inches(3.5), "top": Inches(5.5), "width": Inches(3)},
        "Center": {"left": Inches(3.5), "top": Inches(2.5), "width": Inches(3)}
    }
    
    img_pos = positions.get(image_position, positions["Right Side"])
    
    if show_progress:
        progress_bar = st.progress(0)
        status_text = st.empty()
    
    for idx, slide_data in enumerate(slides_content):
        if show_progress:
            status_text.text(f"Creating slide {idx + 1}/{len(slides_content)}...")
            progress_bar.progress((idx + 1) / len(slides_content))
        
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color_scheme["bg"]
        
        if logo_data:
            try:
                logo_stream = io.BytesIO(logo_data)
                slide.shapes.add_picture(logo_stream, Inches(9), Inches(0.2), width=Inches(0.8))
            except:
                pass
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8.5), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_data["title"]
        title_frame.paragraphs[0].font.size = Pt(36 if idx == 0 else 28)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = color_scheme["accent"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER if idx == 0 else PP_ALIGN.LEFT
        
        if idx > 0 and slide_data.get("bullets"):
            bullet_width = Inches(5.5) if image_mode == "With Images" else Inches(9)
            bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), bullet_width, Inches(4.5))
            text_frame = bullet_box.text_frame
            text_frame.word_wrap = True
            
            for bullet in slide_data["bullets"]:
                p = text_frame.add_paragraph()
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
                p.font.color.rgb = color_scheme["text"]
                p.space_after = Pt(12)
        
        if idx == 0:
            subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = f"{category} Presentation | {audience}"
            subtitle_frame.paragraphs[0].font.size = Pt(20)
            subtitle_frame.paragraphs[0].font.color.rgb = color_scheme["text"]
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]
        
        if idx > 0 and image_mode == "With Images":
            image_prompt = slide_data.get("image_prompt", "")
            
            image_data = get_topic_relevant_image(
                main_topic=topic,
                slide_title=slide_data["title"],
                image_prompt=image_prompt,
                google_api_key=google_api_key,
                google_cx=google_cx,
                use_unsplash=use_unsplash,
                use_pexels=use_pexels,
                pexels_key=pexels_key
            )
            
            if image_data:
                try:
                    image_stream = io.BytesIO(image_data)
                    slide.shapes.add_picture(
                        image_stream, 
                        img_pos["left"], 
                        img_pos["top"], 
                        width=img_pos["width"]
                    )
                except:
                    pass
            
            time.sleep(0.3)
    
    if show_progress:
        progress_bar.progress(1.0)
        status_text.text("✅ Presentation created!")
    
    return prs

# ============================================================================
# LOGIN PAGE
# ============================================================================

def show_login_page():
    """Login page"""
    st.markdown("""
        <style>
        .login-container {
            max-width: 400px;
            margin: 100px auto;
            padding: 40px;
            background: white;
            border-radius: 10px;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
        }
        .login-header {
            text-align: center;
            color: #1f77b4;
            margin-bottom: 30px;
        }
        </style>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    
    with col2:
        st.markdown('<div class="login-container">', unsafe_allow_html=True)
        st.markdown('<div style="font-size: 60px; text-align: center; margin-bottom: 20px;">📊</div>', unsafe_allow_html=True)
        st.markdown('<h1 class="login-header">AI PowerPoint Generator</h1>', unsafe_allow_html=True)
        
        st.markdown("### 🔐 Login")
        
        username = st.text_input("Username", key="login_username")
        password = st.text_input("Password", type="password", key="login_password")
        
        col_btn1, col_btn2 = st.columns(2)
        
        with col_btn1:
            if st.button("🔓 Login", use_container_width=True, type="primary"):
                if username and password:
                    user = verify_user(username, password)
                    if user:
                        st.session_state.logged_in = True
                        st.session_state.user = user
                        log_usage(user['id'], 'login')
                        st.success(f"✅ Welcome, {username}!")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("❌ Invalid credentials")
                else:
                    st.warning("⚠️ Enter username & password")
        
        with col_btn2:
            if st.button("🔑 Demo", use_container_width=True):
                st.info("**Admin:**\nUsername: `admin`\nPassword: `admin123`")
        
        st.markdown('</div>', unsafe_allow_html=True)

# ============================================================================
# INITIALIZE
# ============================================================================

init_database()

# Initialize session state
if 'logged_in' not in st.session_state:
    st.session_state.logged_in = False
if 'user' not in st.session_state:
    st.session_state.user = None
if 'generation_count' not in st.session_state:
    st.session_state.generation_count = 0
if 'total_slides' not in st.session_state:
    st.session_state.total_slides = 0
if 'slides_content' not in st.session_state:
    st.session_state.slides_content = None
if 'google_searches_used' not in st.session_state:
    st.session_state.google_searches_used = 0
if 'templates' not in st.session_state:
    st.session_state.templates = {}
if 'selected_template' not in st.session_state:
    st.session_state.selected_template = None
if 'generation_history' not in st.session_state:
    st.session_state.generation_history = []

# Check login
if not st.session_state.logged_in:
    show_login_page()
    st.stop()

# ============================================================================
# PROFESSIONAL CSS
# ============================================================================

st.markdown("""
<style>
    .main {
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
    }
    
    .main-header {
        font-size: 2.8rem;
        font-weight: 800;
        background: linear-gradient(120deg, #1f77b4, #667eea, #764ba2);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        text-align: center;
        margin-bottom: 0.5rem;
    }
    
    .sub-header {
        text-align: center;
        color: #666;
        font-size: 1.1rem;
        margin-bottom: 2rem;
    }
    
    .dashboard-card {
        background: white;
        padding: 1.5rem;
        border-radius: 15px;
        box-shadow: 0 4px 20px rgba(0,0,0,0.08);
        border: 1px solid rgba(0,0,0,0.05);
        transition: transform 0.3s ease;
    }
    
    .dashboard-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 30px rgba(0,0,0,0.12);
    }
    
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1.2rem;
        border-radius: 12px;
        color: white;
        text-align: center;
        box-shadow: 0 4px 15px rgba(102,126,234,0.4);
    }
    
    .metric-value {
        font-size: 2rem;
        font-weight: 700;
        margin: 0.5rem 0;
    }
    
    .metric-label {
        font-size: 0.85rem;
        opacity: 0.9;
    }
    
    .form-section {
        background: white;
        padding: 1.5rem;
        border-radius: 15px;
        margin: 1rem 0;
        border-left: 5px solid #667eea;
    }
    
    .download-section {
        background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%);
        padding: 2rem;
        border-radius: 20px;
        margin: 2rem 0;
        text-align: center;
        color: white;
        box-shadow: 0 10px 30px rgba(17,153,142,0.3);
    }
    
    .user-info {
        background: #f0f8ff;
        padding: 15px;
        border-radius: 8px;
        margin-bottom: 20px;
    }
    
    .online-indicator {
        display: inline-block;
        width: 10px;
        height: 10px;
        background: #4caf50;
        border-radius: 50%;
        margin-right: 5px;
        animation: pulse 2s infinite;
    }
    
    @keyframes pulse {
        0%, 100% { opacity: 1; }
        50% { opacity: 0.5; }
    }
    
    .activity-card {
        background: #f5f5f5;
        padding: 12px;
        border-radius: 6px;
        margin: 8px 0;
        border-left: 4px solid #1f77b4;
    }
    
    .history-item {
        background: #f8f9fa;
        padding: 1rem;
        border-radius: 10px;
        margin: 0.5rem 0;
        border-left: 4px solid #667eea;
    }
</style>
""", unsafe_allow_html=True)

# ============================================================================
# SIDEBAR (Common for both Admin & User)
# ============================================================================

with st.sidebar:
    user_stats = get_user_stats(st.session_state.user['id'])
    st.markdown(f"""
    <div class='user-info'>
        <h3>👤 {st.session_state.user['username']}</h3>
        <p>Role: <b>{st.session_state.user['role'].upper()}</b></p>
        <hr>
        <p>📊 Presentations: <b>{user_stats['total_presentations']}</b></p>
        <p>📄 Slides: <b>{user_stats['total_slides']}</b></p>
        <p>🔑 Logins: <b>{user_stats['total_logins']}</b></p>
    </div>
    """, unsafe_allow_html=True)
    
    if st.button("🚪 Logout", use_container_width=True):
        logout_user(st.session_state.user['id'])
        log_usage(st.session_state.user['id'], 'logout')
        st.session_state.logged_in = False
        st.session_state.user = None
        st.rerun()
    
    st.markdown("---")
    
    # API Configuration (for PPT generation - both admin & user)
    st.markdown("### ⚙️ Configuration")
    
    with st.expander("🔑 API Keys", expanded=True):
        claude_api_key = st.text_input("OpenRouter API Key", type="password")
        
        model_choice = st.selectbox(
            "AI Model",
            [
                "Free Model (Google Gemini Flash)",
                "Free Model (Meta Llama 3.2)",
                "Free Model (Mistral 7B)",
                "Groq (Llama 3.3 70B) - FREE & FAST",
                "Groq (Mixtral 8x7B) - FREE",
                "Grok-4 Latest (xAI)",
                "Grok-3 (xAI)",
                "Grok-2 (xAI)",
                "Claude 3.5 Sonnet (Paid)"
            ]
        )
        
        groq_api_key = None
        if "Groq" in model_choice:
            groq_api_key = st.text_input("Groq API Key (FREE)", type="password", key="groq_key")
            if groq_api_key:
                st.success("✅ Groq configured!")
        
        grok_api_key = None
        if "Grok" in model_choice:
            grok_api_key = st.text_input("Grok/xAI API Key", type="password", key="grok_key")
            if grok_api_key:
                st.success("✅ Grok configured!")
    
    with st.expander("🖼️ Image Configuration"):
        google_api_key = st.text_input("Google API Key", type="password")
        google_cx = st.text_input("Google CX ID", placeholder="6386765a3a8ed49a9")
        
        if google_api_key and google_cx:
            st.success("✅ Google configured!")
        
        use_unsplash_fallback = st.checkbox("Unsplash", value=True)
        use_pexels_fallback = st.checkbox("Pexels", value=False)
        
        if use_pexels_fallback:
            pexels_api_key = st.text_input("Pexels API Key", type="password")
        else:
            pexels_api_key = None
    
    with st.expander("🏢 Branding"):
        logo_file = st.file_uploader("Company Logo", type=["png", "jpg", "jpeg"])
        logo_data = None
        if logo_file:
            logo_data = logo_file.read()
            st.success("✅ Logo uploaded!")
    
    st.markdown("---")
    
    # Dashboard Metrics
    st.markdown("### 📊 Your Stats")
    col1, col2 = st.columns(2)
    with col1:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-label">Presentations</div>
            <div class="metric-value">{user_stats['total_presentations']}</div>
        </div>
        """, unsafe_allow_html=True)
    with col2:
        st.markdown(f"""
        <div class="metric-card">
            <div class="metric-label">Total Slides</div>
            <div class="metric-value">{user_stats['total_slides']}</div>
        </div>
        """, unsafe_allow_html=True)

# ============================================================================
# ADMIN DASHBOARD
# ============================================================================

if st.session_state.user['role'] == 'admin':
    
    st.markdown('<div class="main-header">👑 Admin Dashboard</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Complete System Management & PPT Generator</div>', unsafe_allow_html=True)
    
    # System Stats
    sys_stats = get_system_stats()
    
    col1, col2, col3, col4, col5 = st.columns(5)
    with col1:
        st.metric("👥 Total Users", sys_stats['total_users'])
    with col2:
        st.metric("🟢 Online Now", sys_stats['currently_online'])
    with col3:
        st.metric("📊 Presentations", sys_stats['total_presentations'])
    with col4:
        st.metric("📄 Slides", sys_stats['total_slides'])
    with col5:
        st.metric("🕒 Today Logins", sys_stats['today_logins'])
    
    st.markdown("---")
    
    # Main Admin Tabs
    admin_main_tab1, admin_main_tab2 = st.tabs([
        "📊 PPT Generator Dashboard",
        "👑 User Management"
    ])
    
    # ========================================================================
    # ADMIN TAB 1: FULL PPT GENERATOR DASHBOARD
    # ========================================================================
    
    with admin_main_tab1:
        st.markdown("## 🚀 AI PowerPoint Generator")
        
        tab1, tab2, tab3, tab4, tab5 = st.tabs([
            "📝 Create", 
            "📁 Templates", 
            "📊 Bulk Generate",
            "📜 History",
            "⚙️ Settings"
        ])
        
        # CREATE TAB
        with tab1:
            st.markdown("### 🚀 Quick Start with Templates")
            
            preset_templates = get_preset_templates()
            cols = st.columns(3)
            selected_preset = None
            
            for idx, (key, template) in enumerate(preset_templates.items()):
                with cols[idx % 3]:
                    if st.button(
                        f"{template['name']}\n{template['description']}", 
                        key=f"admin_preset_{key}",
                        use_container_width=True
                    ):
                        selected_preset = template
                        st.session_state.selected_template = template
            
            st.markdown("---")
            
            col1, col2 = st.columns([1, 1])
            
            with col1:
                st.markdown('<div class="form-section">', unsafe_allow_html=True)
                st.markdown("📝 Content Details", unsafe_allow_html=True)
                
                topic = st.text_input("Topic *", placeholder="e.g., AI in Healthcare", key="admin_topic")
                
                if st.session_state.selected_template:
                    t = st.session_state.selected_template
                    default_category = t.get('category', 'Business')
                    default_slides = t.get('slide_count', 6)
                    default_tone = t.get('tone', 'Formal')
                    default_audience = t.get('audience', 'Corporate')
                    default_theme = t.get('theme', 'Corporate Blue')
                    default_image_mode = t.get('image_mode', 'With Images')
                    default_language = t.get('language', 'English')
                else:
                    default_category = 'Business'
                    default_slides = 6
                    default_tone = 'Formal'
                    default_audience = 'Corporate'
                    default_theme = 'Corporate Blue'
                    default_image_mode = 'With Images'
                    default_language = 'English'
                
                categories = ["Business", "Pitch", "Marketing", "Technical", "Academic", "Training", "Sales"]
                category = st.selectbox(
                    "Category *", 
                    categories,
                    index=categories.index(default_category) if default_category in categories else 0,
                    key="admin_category"
                )
                
                col1_1, col1_2 = st.columns(2)
                with col1_1:
                    slide_count = st.number_input("Slides *", min_value=3, max_value=20, value=default_slides, key="admin_slides")
                with col1_2:
                    languages = ["English", "Hindi (हिंदी)", "Spanish", "French", "German"]
                    language = st.selectbox("Language", languages, key="admin_lang")
                
                tones = ["Formal", "Neutral", "Inspirational", "Educational", "Persuasive"]
                tone = st.selectbox("Tone *", tones, key="admin_tone")
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            with col2:
                st.markdown('<div class="form-section">', unsafe_allow_html=True)
                st.markdown("🎨 Design & Style", unsafe_allow_html=True)
                
                audiences = ["Investors", "Students", "Corporate", "Clients", "Managers"]
                audience = st.selectbox("Target Audience *", audiences, key="admin_audience")
                
                themes_list = ["Corporate Blue", "Gradient Modern", "Minimal Dark", "Pastel Soft", "Professional Green", "Elegant Purple"]
                theme = st.selectbox("Visual Theme *", themes_list, key="admin_theme")
                
                image_modes = ["With Images", "No Images"]
                image_mode = st.selectbox("Image Mode *", image_modes, key="admin_imgmode")
                
                if image_mode == "With Images":
                    image_position = st.selectbox("Image Position", ["Right Side", "Left Side", "Top Right Corner", "Bottom", "Center"], key="admin_imgpos")
                else:
                    image_position = "Right Side"
                
                st.markdown('</div>', unsafe_allow_html=True)
            
            with st.expander("➕ Additional Options"):
                key_points = st.text_area("Key Points", placeholder="- Point 1\n- Point 2", key="admin_keypoints")
                export_format = st.selectbox("Export Format", ["PowerPoint (.pptx)", "PowerPoint + PDF", "Google Slides (JSON)"], key="admin_export")
            
            st.markdown("---")
            
            col_btn1, col_btn2 = st.columns([2, 1])
            
            with col_btn1:
                generate_button = st.button("🚀 Generate Presentation", use_container_width=True, type="primary", key="admin_generate")
            
            with col_btn2:
                save_as_template = st.button("💾 Save as Template", use_container_width=True, key="admin_save_template")
            
            if generate_button and topic:
                has_valid_api = False
                if "Groq" in model_choice and groq_api_key:
                    has_valid_api = True
                elif "Grok" in model_choice and grok_api_key:
                    has_valid_api = True
                elif claude_api_key:
                    has_valid_api = True
                
                if has_valid_api:
                    with st.spinner("🤖 Generating..."):
                        slides_content = generate_content_with_retry(
                            claude_api_key, topic, category, slide_count, 
                            tone, audience, key_points, model_choice, language,
                            grok_api_key=grok_api_key,
                            groq_api_key=groq_api_key
                        )
                        
                        if slides_content:
                            log_usage(st.session_state.user['id'], 'generate_presentation', topic, len(slides_content))
                            
                            prs = create_powerpoint(
                                slides_content, theme, image_mode,
                                google_api_key if 'google_api_key' in locals() else "",
                                google_cx if 'google_cx' in locals() else "",
                                use_unsplash_fallback, use_pexels_fallback, 
                                pexels_api_key if use_pexels_fallback and 'pexels_api_key' in locals() else None,
                                category, audience, topic, 
                                image_position, logo_data
                            )
                            
                            pptx_io = io.BytesIO()
                            prs.save(pptx_io)
                            pptx_io.seek(0)
                            
                            st.markdown("""
                            <div class="download-section">
                                <h2>🎉 Presentation Ready!</h2>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            st.download_button(
                                "📥 DOWNLOAD POWERPOINT",
                                pptx_io.getvalue(),
                                f"{topic.replace(' ', '_')}.pptx",
                                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True,
                                type="primary"
                            )
                            
                            with st.expander("📄 Preview"):
                                for idx, slide in enumerate(slides_content):
                                    st.markdown(f"### Slide {idx + 1}: {slide['title']}")
                                    if slide.get('bullets'):
                                        for bullet in slide['bullets']:
                                            st.markdown(f"• {bullet}")
                                    st.markdown("---")
        
        # TEMPLATES TAB
        with tab2:
            st.markdown("### 📁 Template Manager")
            if st.session_state.templates:
                for temp_id, template in st.session_state.templates.items():
                    st.markdown(f"**{template['name']}** - {template['category']} | {template['slide_count']} slides")
                    if st.button("Use", key=f"admin_use_{temp_id}"):
                        st.session_state.selected_template = template
            else:
                st.info("No templates saved yet")
        
        # BULK GENERATE TAB
        with tab3:
            st.markdown("### 📊 Bulk Generate")
            st.info("Upload CSV with multiple presentation topics")
            uploaded_csv = st.file_uploader("Upload CSV", type=['csv'], key="admin_bulk_csv")
        
        # HISTORY TAB
        with tab4:
            st.markdown("### 📜 History")
            my_activities = get_user_activity_details(st.session_state.user['id'])
            if my_activities:
                for activity in my_activities[:10]:
                    action, topic, slides_count, timestamp = activity
                    if action == 'generate_presentation':
                        st.markdown(f"📊 {topic} | {slides_count} slides | {timestamp}")
        
        # SETTINGS TAB
        with tab5:
            st.markdown("### ⚙️ Settings")
            st.info("Settings panel")
    
    # ========================================================================
    # ADMIN TAB 2: USER MANAGEMENT
    # ========================================================================
    
    with admin_main_tab2:
        st.markdown("## 👑 User Management")
        
        user_tab1, user_tab2, user_tab3, user_tab4 = st.tabs([
            "🟢 Live Dashboard",
            "➕ Create User",
            "👥 Manage Users",
            "📊 Activity Log"
        ])
        
        with user_tab1:
            st.markdown("### 🟢 Live User Activity")
            
            if st.button("🔄 Refresh", key="admin_refresh"):
                st.rerun()
            
            active_users = get_currently_logged_in_users()
            
            if active_users:
                st.success(f"**{len(active_users)} user(s) online**")
                
                for user in active_users:
                    user_id, username, email, login_time, role = user
                    user_activity_stats = get_user_stats(user_id)
                    
                    st.markdown(f"""
<div style='background:#e8f5e9;padding:15px;border-radius:8px;margin:10px 0;'>
    <span class='online-indicator'></span>
    <b>{username}</b> ({role})<br>
    <small>📧 {email if email else 'N/A'} | 🕒 {login_time}</small><br>
    <small>📊 {user_activity_stats['total_presentations']} ppts | 📄 {user_activity_stats['total_slides']} slides</small>
</div>
                    """, unsafe_allow_html=True)
            else:
                st.warning("No users online")
        
        with user_tab2:
            st.markdown("### ➕ Create User")
            
            with st.form("admin_create_user"):
                col_a, col_b = st.columns(2)
                with col_a:
                    new_username = st.text_input("Username *")
                    new_email = st.text_input("Email")
                with col_b:
                    new_password = st.text_input("Password *", type="password")
                    confirm_password = st.text_input("Confirm *", type="password")
                
                submitted = st.form_submit_button("✅ Create User", type="primary")
                
                if submitted:
                    if new_username and new_password == confirm_password and len(new_password) >= 6:
                        if create_user_by_admin(new_username, new_password, new_email):
                            st.success(f"""
✅ User Created!

Username: `{new_username}`
Password: `{new_password}`
                            """)
                        else:
                            st.error("Username exists!")
        
        with user_tab3:
            st.markdown("### 👥 All Users")
            
            users = get_all_users()
            user_data = []
            for user in users:
                user_data.append({
                    'ID': user[0],
                    'Username': user[1],
                    'Email': user[2] if user[2] else 'N/A',
                    'Created': user[3],
                    'Active': '✅' if user[5] else '❌',
                    'Role': user[6]
                })
            
            df_users = pd.DataFrame(user_data)
            st.dataframe(df_users, use_container_width=True)
            
            st.markdown("---")
            col_m1, col_m2, col_m3 = st.columns(3)
            with col_m1:
                user_id_action = st.number_input("User ID", min_value=1, step=1)
            with col_m2:
                action_type = st.selectbox("Action", ["Enable", "Disable", "Delete"])
            with col_m3:
                st.write("")
                if st.button("▶️ Execute", type="primary"):
                    if user_id_action != 1:
                        if action_type == "Enable":
                            toggle_user_status(user_id_action, 1)
                            st.success("Enabled!")
                            time.sleep(1)
                            st.rerun()
                        elif action_type == "Disable":
                            toggle_user_status(user_id_action, 0)
                            st.warning("Disabled!")
                            time.sleep(1)
                            st.rerun()
                        elif action_type == "Delete":
                            delete_user(user_id_action)
                            st.error("Deleted!")
                            time.sleep(1)
                            st.rerun()
        
        with user_tab4:
            st.markdown("### 📊 Activity Log")
            
            all_activities = get_all_user_activities()
            
            if all_activities:
                activity_records = []
                for activity in all_activities:
                    username, action, topic, slides_count, timestamp = activity
                    activity_records.append({
                        'Username': username,
                        'Action': action,
                        'Topic': topic if topic else '-',
                        'Slides': slides_count if slides_count else '-',
                        'Timestamp': timestamp
                    })
                
                df_activities = pd.DataFrame(activity_records)
                st.dataframe(df_activities, use_container_width=True, height=600)

# ============================================================================
# USER DASHBOARD (Full PPT Generator)
# ============================================================================

else:
    
    st.markdown('<div class="main-header">📊 AI PowerPoint Generator Pro</div>', unsafe_allow_html=True)
    st.markdown('<div class="sub-header">Create stunning presentations with AI-powered content</div>', unsafe_allow_html=True)
    
    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "📝 Create", 
        "📁 Templates", 
        "📊 Bulk Generate",
        "📜 History",
        "⚙️ Settings"
    ])
    
    # CREATE TAB
    with tab1:
        st.markdown("### 🚀 Quick Start with Templates")
        
        preset_templates = get_preset_templates()
        cols = st.columns(3)
        
        for idx, (key, template) in enumerate(preset_templates.items()):
            with cols[idx % 3]:
                if st.button(
                    f"{template['name']}\n{template['description']}", 
                    key=f"user_preset_{key}",
                    use_container_width=True
                ):
                    st.session_state.selected_template = template
        
        st.markdown("---")
        
        col1, col2 = st.columns([1, 1])
        
        with col1:
            st.markdown('<div class="form-section">', unsafe_allow_html=True)
            st.markdown("📝 Content Details", unsafe_allow_html=True)
            
            topic = st.text_input("Topic *", placeholder="e.g., AI in Healthcare", key="user_topic")
            
            if st.session_state.selected_template:
                t = st.session_state.selected_template
                default_category = t.get('category', 'Business')
                default_slides = t.get('slide_count', 6)
                default_tone = t.get('tone', 'Formal')
                default_audience = t.get('audience', 'Corporate')
                default_theme = t.get('theme', 'Corporate Blue')
                default_image_mode = t.get('image_mode', 'With Images')
                default_language = t.get('language', 'English')
            else:
                default_category = 'Business'
                default_slides = 6
                default_tone = 'Formal'
                default_audience = 'Corporate'
                default_theme = 'Corporate Blue'
                default_image_mode = 'With Images'
                default_language = 'English'
            
            categories = ["Business", "Pitch", "Marketing", "Technical", "Academic", "Training", "Sales"]
            category = st.selectbox("Category *", categories, key="user_category")
            
            col1_1, col1_2 = st.columns(2)
            with col1_1:
                slide_count = st.number_input("Slides *", min_value=3, max_value=20, value=default_slides, key="user_slides")
            with col1_2:
                languages = ["English", "Hindi (हिंदी)", "Spanish", "French", "German"]
                language = st.selectbox("Language", languages, key="user_lang")
            
            tones = ["Formal", "Neutral", "Inspirational", "Educational", "Persuasive"]
            tone = st.selectbox("Tone *", tones, key="user_tone")
            
            st.markdown('</div>', unsafe_allow_html=True)
        
        with col2:
            st.markdown('<div class="form-section">', unsafe_allow_html=True)
            st.markdown("🎨 Design & Style", unsafe_allow_html=True)
            
            audiences = ["Investors", "Students", "Corporate", "Clients", "Managers"]
            audience = st.selectbox("Target Audience *", audiences, key="user_audience")
            
            themes_list = ["Corporate Blue", "Gradient Modern", "Minimal Dark", "Pastel Soft", "Professional Green", "Elegant Purple"]
            theme = st.selectbox("Visual Theme *", themes_list, key="user_theme")
            
            image_modes = ["With Images", "No Images"]
            image_mode = st.selectbox("Image Mode *", image_modes, key="user_imgmode")
            
            if image_mode == "With Images":
                image_position = st.selectbox("Image Position", ["Right Side", "Left Side", "Top Right Corner", "Bottom", "Center"], key="user_imgpos")
            else:
                image_position = "Right Side"
            
            st.markdown('</div>', unsafe_allow_html=True)
        
        with st.expander("➕ Additional Options"):
            key_points = st.text_area("Key Points", placeholder="- Point 1\n- Point 2", key="user_keypoints")
            export_format = st.selectbox("Export Format", ["PowerPoint (.pptx)", "PowerPoint + PDF", "Google Slides (JSON)"], key="user_export")
        
        st.markdown("---")
        
        if st.button("🚀 Generate Presentation", use_container_width=True, type="primary", key="user_generate"):
            if topic:
                has_valid_api = False
                if "Groq" in model_choice and groq_api_key:
                    has_valid_api = True
                elif "Grok" in model_choice and grok_api_key:
                    has_valid_api = True
                elif claude_api_key:
                    has_valid_api = True
                
                if has_valid_api:
                    with st.spinner("🤖 Generating your presentation..."):
                        slides_content = generate_content_with_retry(
                            claude_api_key, topic, category, slide_count, 
                            tone, audience, key_points, model_choice, language,
                            grok_api_key=grok_api_key,
                            groq_api_key=groq_api_key
                        )
                        
                        if slides_content:
                            log_usage(st.session_state.user['id'], 'generate_presentation', topic, len(slides_content))
                            
                            prs = create_powerpoint(
                                slides_content, theme, image_mode,
                                google_api_key if 'google_api_key' in locals() else "",
                                google_cx if 'google_cx' in locals() else "",
                                use_unsplash_fallback, use_pexels_fallback, 
                                pexels_api_key if use_pexels_fallback and 'pexels_api_key' in locals() else None,
                                category, audience, topic, 
                                image_position, logo_data
                            )
                            
                            pptx_io = io.BytesIO()
                            prs.save(pptx_io)
                            pptx_io.seek(0)
                            
                            st.markdown("""
                            <div class="download-section">
                                <h2>🎉 Your Presentation is Ready!</h2>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            st.download_button(
                                "📥 DOWNLOAD POWERPOINT",
                                pptx_io.getvalue(),
                                f"{topic.replace(' ', '_')}.pptx",
                                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                use_container_width=True,
                                type="primary"
                            )
                            
                            st.balloons()
                            
                            with st.expander("📄 Preview Slides", expanded=True):
                                for idx, slide in enumerate(slides_content):
                                    st.markdown(f"### Slide {idx + 1}: {slide['title']}")
                                    if slide.get('bullets'):
                                        for bullet in slide['bullets']:
                                            st.markdown(f"• {bullet}")
                                    st.markdown("---")
                            
                            with st.expander("🎓 AI Coach"):
                                issues, suggestions, score = analyze_presentation(slides_content)
                                st.metric("Quality Score", f"{score}/100")
                                if suggestions:
                                    for suggestion in suggestions:
                                        st.write(f"• {suggestion}")
                else:
                    st.error("⚠️ Please configure API keys in sidebar")
    
    # TEMPLATES TAB
    with tab2:
        st.markdown("### 📁 Template Manager")
        if st.session_state.templates:
            for temp_id, template in st.session_state.templates.items():
                st.markdown(f"**{template['name']}** - {template['category']} | {template['slide_count']} slides")
                if st.button("Use", key=f"user_use_{temp_id}"):
                    st.session_state.selected_template = template
        else:
            st.info("No templates saved yet")
    
    # BULK GENERATE TAB
    with tab3:
        st.markdown("### 📊 Bulk Generate")
        st.info("Upload CSV with multiple presentation topics")
    
    # HISTORY TAB
    with tab4:
        st.markdown("### 📜 Your History")
        my_activities = get_user_activity_details(st.session_state.user['id'])
        if my_activities:
            for activity in my_activities[:10]:
                action, topic, slides_count, timestamp = activity
                if action == 'generate_presentation':
                    st.markdown(f"""
<div class='activity-card'>
    <b>📊 {topic}</b><br>
    <small>📄 {slides_count} slides | 🕒 {timestamp}</small>
</div>
                    """, unsafe_allow_html=True)
        else:
            st.info("No history yet")
    
    # SETTINGS TAB
    with tab5:
        st.markdown("### ⚙️ Settings")
        st.info("Settings panel")

# Footer
st.markdown("---")
st.markdown(f"""
<div style='text-align: center; color: #666;'>
    <p>🔐 Logged in as: <b>{st.session_state.user['username']}</b> ({st.session_state.user['role']}) | {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}</p>
    <p>✨ AI PowerPoint Generator Pro | Version 3.0</p>
</div>
""", unsafe_allow_html=True)
